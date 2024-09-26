from docx import Document
import pypdf as PyPDF2
from pptx import Presentation
import arabic_reshaper
from bidi.algorithm import get_display
import re


import os
import streamlit as st

FILE_TYPE_DOCX = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
FILE_TYPE_PDF = "application/pdf"
FILE_TYPE_PPTX = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
FILE_TYPE_TEXT = "text/plain"
text_file_types = [
    FILE_TYPE_DOCX, FILE_TYPE_PDF, FILE_TYPE_PPTX, FILE_TYPE_TEXT
]


class UploadedFileWrapper:
    def __init__(self, file_path):
        self.file = open(file_path, "rb")
        self.name = os.path.basename(file_path)
        self.type = self.get_file_type(file_path)
        self.size = os.path.getsize(file_path)

    def get_file_type(self, file_path):
        if file_path.endswith(".txt"):
            return FILE_TYPE_TEXT
        elif file_path.endswith(".pdf"):
            return FILE_TYPE_PDF
        elif file_path.endswith(".docx"):
            return FILE_TYPE_DOCX
        else:
            return ""

    def read(self):
        return self.file.read()

    def __del__(self):
        self.file.close()


def parse_file(file):
    if file.type == FILE_TYPE_DOCX:
        return text_from_docx(file)
    elif file.type == FILE_TYPE_PDF:
        return text_from_pdf(file)
    elif file.type == FILE_TYPE_TEXT:
        return file.read().decode("utf-8")
    elif file.type == FILE_TYPE_PPTX:
        return text_from_pptx(file)
    return


def text_from_docx(docx_file):
    doc = Document(docx_file)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():  # Only add non-empty strings
            full_text.append(para.text)
    return "\n".join(full_text)


def is_arabic(text):
    arabic_pattern = re.compile("[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF]+")
    return bool(arabic_pattern.search(text))


def text_from_pdf(file):
    if isinstance(file, UploadedFileWrapper):
        reader = PyPDF2.PdfReader(file.file)
    else: 
        reader = PyPDF2.PdfReader(file)
    num_pages = len(reader.pages)

    full_text = ""
    for page_num in range(num_pages):
        page = reader.pages[page_num]
        text = page.extract_text()
        if is_arabic(text):
            # Reshape Arabic text
            reshaped_text = arabic_reshaper.reshape(text)
            # Correct text direction
            bidi_text = get_display(reshaped_text)
            full_text += bidi_text + "\n\n"
        else:
            # For non-Arabic (presumably English) text, just add it as is
            full_text += text + "\n\n"
    return full_text


def process_text(text):
    if is_arabic(text):
        # Reshape Arabic text
        reshaped_text = arabic_reshaper.reshape(text)
        # Correct text direction
        return get_display(reshaped_text)
    return text  # Return as-is for non-Arabic text


def text_from_pptx(file_path):
    prs = Presentation(file_path)
    full_text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                processed_text = process_text(shape.text)
                if processed_text.strip():  # Only add non-empty strings
                    slide_text.append(processed_text)

        if slide_text:
            full_text.append("\n".join(slide_text))
    return "\n".join(full_text)
